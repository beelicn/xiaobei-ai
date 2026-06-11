import os
import tempfile
from openai import OpenAI
from docx import Document
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io
from PyPDF2 import PdfReader
import streamlit as st
from pdf2docx import Converter

# ==============================================================================
# 🌐 【多语言配置区】中英文语言包
# ==============================================================================
LANG_PACK = {
    "zh": {
        # 全局通用
        "page_title": "小倍AI助手",
        "main_title": "🙂 小倍AI助手",
        "warning": "警告",
        "error": "错误",
        "success": "完成",
        "start": "开始",
        "download": "下载",
        "upload": "上传",
        "preview": "预览",
        "generating": "生成中...",
        "processing": "处理中...",
        # 侧边栏
        "sidebar_title": "功能导航",
        "select_func": "选择功能",
        "sidebar_footer": "✅修复了一些已知问题ㅤㅤㅤ☺️版本：测试版 ㅤㅤㅤ©️beelicn.com",
        "lang_select": "选择语言",
        # 功能菜单（label+副标题）— 7个原有功能
        "menu_search_label": "💻 全网报告搜索",
        "menu_search_sub": "🤔 小倍正在搜索报告中",
        "menu_summary_label": "💡 文档总结/数据提取",
        "menu_summary_sub": "🤨 小倍正在分析你的文档",
        "menu_generate_label": "📝 行业报告生成",
        "menu_generate_sub": "🤯 小倍正在生成你的报告",
        "menu_compare_label": "📈 多文档竞品/赛道对比分析",
        "menu_compare_sub": "😁 小倍正在对比中",
        "menu_rewrite_label": "✏️ 仿照模板改写文档",
        "menu_rewrite_sub": "🧐 小倍正在改写你的文档",
        "menu_translate_label": "🌐 商务文档翻译",
        "menu_translate_sub": "😏 小倍翻译中",
        "menu_pdf2word_label": "💾 无损PDF转Word",
        "menu_pdf2word_sub": "😎 小倍PDF格式转换助手",
        # 新增功能菜单
        "menu_about_label": "📖 关于本产品",
        "menu_about_sub": "产品方法论与设计理念",
        "menu_assessment_label": "🏢 企业AI就绪度评估",
        "menu_assessment_sub": "评估组织的AI准备情况",
        "menu_compliance_label": "🛡️ AI合规与溯源",
        "menu_compliance_sub": "可信度评分与合规机制",
        "menu_dashboard_label": "📊 战略决策仪表盘",
        "menu_dashboard_sub": "决策优先级矩阵与战略建议",
        # 赛道配置
        "track_general": "通用全行业",
        "track_ai": "AI市场研究",
        "track_consulting": "战略咨询",
        "track_risk": "企业风险管理",
        "track_manufacture": "制造业出海欧洲市场",
        "track_esg": "ESG/可持续发展",
        # 功能1：全网报告搜索
        "search_input_tip": "请输入行业/赛道关键词，越详细检索结果越精准",
        "search_btn": "开始检索",
        "search_loading": "正在检索合规行业报告...",
        "search_kw_empty": "请输入检索关键词",
        "search_pub_org": "🏢 发布机构：",
        "search_pub_year": "📅 发布年份：",
        "search_abstract": "📄 核心摘要：",
        "search_credibility": "🔍 可信度评分：",
        # 功能2：文档总结/数据提取
        "summary_mode": "选择分析模式",
        "summary_mode_general": "通用文档总结",
        "summary_mode_indicator": "行研核心指标提取",
        "summary_upload_tipip": "上传TXT/DOCX格式的文档、财报、行业白皮书",
        "summary_upload_tip": "上传TXT/DOCX格式的文档、财报、行业白皮书",
        "summary_analyze_btn": "开始分析",
        "summary_analyze_loading": "正在执行{mode}...",
        "summary_original_preview": "文档原文预览",
        "summary_result_title": "✅ {mode}结果",
        "summary_download_btn": "📍 下载分析结果Word文档",
        "summary_download_filename": "文档{mode}结果.docx",
        # 功能3：行业报告生成
        "generate_track_select": "选择垂直赛道模板",
        "generate_name_input": "输入目标行业/赛道/产品名称",
        "generate_ref_tip": "【可选】上传自有参考资料/报告模板（生成内容优先匹配参考资料的格式与规范）",
        "generate_ref_upload": "上传参考资料TXT/DOCX文档",
        "generate_ref_preview": "预览参考资料内容",
        "generate_btn": "📝 生成行业报告",
        "generate_loading": "正在生成{track}赛道专属行业报告...",
        "generate_name_empty": "请输入目标行业/赛道名称",
        "generate_ref_rule": "【参考资料要求】生成内容必须优先参考以下资料的格式规范、行业定义、数据口径：",
        "generate_report_title": "✉️ {name} | {track}赛道行业报告",
        "generate_download_word": "📍 下载Word版报告",
        "generate_download_ppt": "📊 下载咨询标准PPT版",
        "generate_word_filename": "{name}_{track}_行业报告.docx",
        "generate_ppt_filename": "{name}_{track}_行业报告.pptx",
        "generate_compliance_checkbox": "⚖️ 在报告中加入法律合规风险分析模块",
        # 功能4：多文档对比分析
        "compare_tip": "支持上传2-5份同赛道行业报告、竞品财报、行业白皮书，自动生成战略咨询级对比分析报告",
        "compare_upload_tip": "上传需要对比分析的TXT/DOCX文档",
        "compare_btn": "📈 生成对比分析报告",
        "compare_loading": "正在解析文档并生成对比分析报告...",
        "compare_file_min": "请至少上传2份文档进行对比分析",
        "compare_result_title": "✅ 对比分析报告",
        "compare_download_word": "📍 下载Word版分析报告",
        "compare_download_ppt": "📊 下载咨询标准PPT版",
        "compare_word_filename": "赛道竞品对比分析报告.docx",
        "compare_ppt_filename": "赛道竞品对比分析报告.pptx",
        "compare_priority_label": "📊 包含战略建议优先级矩阵",
        # 功能5：仿照模板改写文档
        "rewrite_flow": "流程：上传模板文档 → 上传待改写文档 → 一键改写 → 在线预览 → 下载双格式文件",
        "rewrite_template_upload": "1. 上传模板文档",
        "rewrite_content_upload": "2. 上传待改写文档",
        "rewrite_template_preview": "预览模板内容",
        "rewrite_content_preview": "预览待改写内容",
        "rewrite_btn": "✏️ 开始改写",
        "rewrite_loading": "正在按模板风格改写文档...",
        "rewrite_file_empty": "请先上传模板文档和待改写文档",
        "rewrite_result_title": "✅ 改写结果",
        "rewrite_download_word": "📍 下载Word文档",
        "rewrite_download_ppt": "📊 下载PPT版",
        "rewrite_word_filename": "文档改写结果.docx",
        "rewrite_ppt_filename": "文档改写结果.pptx",
        # 功能6：商务文档翻译
        "translate_tip": "支持直接输入文本翻译，或上传TXT/DOCX文档批量翻译，适配商务/咨询正式文档场景",
        "translate_target_lang": "选择目标翻译语言",
        "translate_mode": "翻译模式",
        "translate_mode_text": "直接输入文本",
        "translate_mode_file": "上传文档翻译",
        "translate_textarea_tip": "请输入需要翻译的商务文档内容",
        "translate_upload_tip": "上传需要翻译的TXT/DOCX文档",
        "translate_original_preview": "预览原文内容",
        "translate_btn": "🌐 开始翻译",
        "translate_loading": "正在翻译中，请稍等...",
        "translate_content_empty": "请输入需要翻译的内容，或上传有效文档",
        "translate_result_title": "✅ 翻译结果",
        "translate_download_btn": "📍 下载翻译结果Word文档",
        "translate_download_filename": "商务文档翻译结果.docx",
        # 功能7：PDF转Word
        "pdf2word_tip": "上传PDF → AI智能修复乱换行/乱分段 → 还原整洁排版 → 预览下载双格式文件",
        "pdf2word_upload_tip": "上传PDF文件",
        "pdf2word_loading": "正在提取PDF内容，并AI智能规整排版...",
        "pdf2word_preview_title": "📋 AI规整后内容预览",
        "pdf2word_download_word": "📍 下载无损Word文档",
        "pdf2word_download_ppt": "📊 下载PPT版",
        "pdf2word_word_filename": "PDF转换结果.docx",
        "pdf2word_ppt_filename": "PDF转换结果.pptx",
        # 新增：关于本产品
        "about_title": "关于本产品",
        "about_subtitle": "产品方法论与设计理念",
        "about_problem": "一、问题定义",
        "about_solution": "二、解决方案",
        "about_architecture": "三、技术架构",
        "about_value": "四、商业价值",
        "about_principles": "五、设计原则",
        "about_content": """### 一、问题定义
商业决策者面临信息过载、数据碎片化、专业报告获取成本高的困境。传统行业研究依赖人工搜集与整理，效率低下且容易遗漏关键信息。

### 二、解决方案
小倍AI助手通过大语言模型与专业提示词工程，提供从报告搜索、文档分析、报告生成到竞品对比的一站式AI商业决策支持平台，大幅降低专业研究门槛。

### 三、技术架构
- **大模型层**：火山引擎豆包大模型（doubao-seed-2-0-pro），提供高质量语义理解与生成
- **应用层**：Streamlit前端框架，Python后端
- **文档处理**：python-docx / python-pptx / PyPDF2 / pdf2docx
- **合规机制**：数据来源可信度评分、AI生成内容溯源、合规检查开关

### 四、商业价值
- 将咨询级行业分析成本降低90%以上
- 报告生成时间从天级压缩到分钟级
- 覆盖7大核心商业决策场景
- 支持中英文双语专业输出

### 五、设计原则
- **负责任AI**：所有生成内容可溯源、可验证
- **专业主义**：严格遵循咨询行业报告规范
- **跨学科融合**：融合法律合规、ESG、组织变革等视角
- **产品方法论驱动**：从问题定义到价值交付的完整闭环""",
        # 新增：企业AI就绪度评估
        "assessment_title": "企业AI就绪度评估",
        "assessment_subtitle": "评估您所在组织的AI应用准备情况，获得定制化建议",
        "assessment_company": "公司名称",
        "assessment_industry": "所属行业",
        "assessment_size": "员工规模",
        "assessment_revenue": "年营收（万元人民币）",
        "assessment_existing_ai": "现有AI应用情况",
        "assessment_btn": "🔍 开始评估",
        "assessment_loading": "正在评估贵司AI就绪度...",
        "assessment_result_title": "📋 AI就绪度评估报告",
        "assessment_sizes": ["<100人", "100-500人", "500-2000人", "2000人以上"],
        "assessment_ai_levels": ["尚未使用AI", "初步尝试（如ChatGPT）", "已在部分业务中应用AI工具", "已建立AI团队/平台"],
        "assessment_company_empty": "请输入公司名称",
        "assessment_fields": {
            "digital_level": "数字化水平评估",
            "ai_replace": "AI可替代工作清单",
            "org_advice": "组织架构调整建议",
            "roi_estimate": "ROI预估"
        },
        # 新增：AI合规与溯源
        "compliance_title": "🛡️ AI合规与溯源机制",
        "compliance_subtitle": "负责任AI的三大核心机制",
        "compliance_score_title": "一、数据来源可信度评分体系",
        "compliance_score_desc": "对AI生成内容中引用的数据来源进行可信度分级评分，确保信息的可靠性和可验证性。",
        "compliance_trace_title": "二、AI生成内容溯源链路",
        "compliance_trace_desc": "建立从用户输入→提示词→模型输出→最终呈现的完整溯源链路，每一步均可审计。",
        "compliance_switch_title": "三、合规检查开关",
        "compliance_switch_desc": "用户可在生成报告前选择开启/关闭合规检查项，灵活平衡效率与合规需求。",
        "compliance_source_levels": "可信度分级",
        "compliance_source_gov": "政府/监管机构发布",
        "compliance_source_listed": "上市公司官方披露",
        "compliance_source_consulting": "头部咨询/研究机构",
        "compliance_source_media": "权威财经/科技媒体",
        "compliance_source_unknown": "未标注来源",
        "compliance_trace_flow": "用户输入 → 提示词工程（含合规规则注入） → 豆包大模型推理 → 合规后处理（可信度校验） → 最终输出",
        "compliance_switch_items": [
            "AI生成内容合规检查（自动检测虚构数据、不当引用）",
            "数据来源溯源（要求AI标注每条数据的出处）",
            "偏见与幻觉检测（交叉验证关键事实）"
        ],
        # 新增：战略决策仪表盘（交互式）
        "dashboard_title": "📊 战略决策仪表盘",
        "dashboard_subtitle": "AI驱动的艾森豪威尔矩阵与优先级排序",
        "dashboard_intro": "在下方输入您的项目/方案清单，AI将生成艾森豪威尔矩阵（四象限优先级框架）并输出P0-P3执行优先级。",
        "dashboard_projects_tip": "请输入您的项目/方案（每行一个）",
        "dashboard_projects_placeholder": "例如：\nQ3新品上线活动\n官网改版项目\n竞品定价分析报告\nCRM系统升级评估",
        "dashboard_btn": "📊 生成决策矩阵",
        "dashboard_loading": "AI正在生成您的艾森豪威尔矩阵...",
        "dashboard_result_title": "✅ 战略决策矩阵",
        "dashboard_download_word": "📍 下载Word版矩阵报告",
        "dashboard_download_ppt": "📊 下载PPT版矩阵报告",
        # 新增：AI合规与溯源（交互式）
        "compliance_title": "🛡️ AI合规与溯源",
        "compliance_subtitle": "负责任AI：合规检查 + 溯源报告生成",
        "compliance_check_title": "一、合规检查（粘贴文本进行审查）",
        "compliance_check_tip": "在下方粘贴AI生成内容或商务文档文本，进行合规审查",
        "compliance_check_placeholder": "在此粘贴待审查的文本...",
        "compliance_check_btn": "🔍 执行合规检查",
        "compliance_check_loading": "正在审查文本合规性...",
        "compliance_check_result_title": "✅ 合规检查结果",
        "compliance_trace_title_new": "二、溯源报告生成器",
        "compliance_trace_tip": "输入报告主题，AI将生成带完整来源标注的报告框架",
        "compliance_trace_placeholder": "例如：2024年生成式AI行业报告",
        "compliance_trace_btn": "📋 生成溯源报告",
        "compliance_trace_loading": "正在生成带来源标注的溯源报告...",
        "compliance_trace_result_title": "✅ 溯源报告",
        "compliance_download_word": "📍 下载Word版报告",
        "compliance_download_ppt": "📊 下载PPT版报告",
        # 静态合规参考（置于交互模块下方）
        "compliance_static_title": "三、合规机制参考",
        "compliance_score_title": "数据来源可信度评分体系",
        "compliance_score_desc": "对AI生成内容中引用的数据来源进行可信度分级评分，确保信息的可靠性和可验证性。",
        "compliance_trace_title": "AI生成内容溯源链路",
        "compliance_trace_desc": "建立从用户输入→提示词→模型输出→最终呈现的完整溯源链路，每一步均可审计。",
        "compliance_switch_title": "合规检查开关",
        "compliance_switch_desc": "用户可在生成报告前选择开启/关闭合规检查项，灵活平衡效率与合规需求。",
        "compliance_source_levels": "可信度分级",
        "compliance_source_gov": "政府/监管机构发布",
        "compliance_source_listed": "上市公司官方披露",
        "compliance_source_consulting": "头部咨询/研究机构",
        "compliance_source_media": "权威财经/科技媒体",
        "compliance_source_unknown": "未标注来源",
        "compliance_trace_flow": "用户输入 → 提示词工程（含合规规则注入） → 豆包大模型推理 → 合规后处理（可信度校验） → 最终输出",
        "compliance_switch_items": [
            "AI生成内容合规检查（自动检测虚构数据、不当引用）",
            "数据来源溯源（要求AI标注每条数据的出处）",
            "偏见与幻觉检测（交叉验证关键事实）"
        ],
        # 其他
        "func_not_found": "该功能暂未实现，请检查配置",
        "ppt_title_default": "咨询报告",
        "ppt_footer": "小倍咨询级AI报告助手\n合规生成 | 数据可溯源",
        "ppt_end_page": "报告结束",
        "ai_empty_response": "AI生成内容为空，请重试~",
    },
    "en": {
        # Global General
        "page_title": "Xiaobei AI Assistant",
        "main_title": "🙂 Xiaobei AI Assistant",
        "warning": "Warning",
        "error": "Error",
        "success": "Success",
        "start": "Start",
        "download": "Download",
        "upload": "Upload",
        "preview": "Preview",
        "generating": "Generating...",
        "processing": "Processing...",
        # Sidebar
        "sidebar_title": "Function Navigation",
        "select_func": "Select Function",
        "sidebar_footer": "✅ Fixed known issuesㅤㅤㅤㅤ☺️Version: Beta ㅤㅤ©️beelicn.com",
        "lang_select": "Select Language",
        # Menu Config — 7 original
        "menu_search_label": "💻 Full-web Report Search",
        "menu_search_sub": "🤔 Xiaobei is searching reports",
        "menu_summary_label": "💡 Doc Sum/Data Extraction",
        "menu_summary_sub": "🤨 Xiaobei is analyzing your document",
        "menu_generate_label": "📝 Industry Report Generation",
        "menu_generate_sub": "🤯 Xiaobei is generating your report",
        "menu_compare_label": "📈  Competitor Analysis",
        "menu_compare_sub": "😁 Xiaobei is comparing documents",
        "menu_rewrite_label": "✏️ Template-based Rewrite",
        "menu_rewrite_sub": "🧐 Xiaobei is rewriting your document",
        "menu_translate_label": "🌐 Business Translation",
        "menu_translate_sub": "😏 Xiaobei is translating",
        "menu_pdf2word_label": "💾 Lossless PDF to Word",
        "menu_pdf2word_sub": "😎 Xiaobei PDF Converter",
        # New menu items
        "menu_about_label": "📖 About This Product",
        "menu_about_sub": "Product Methodology & Design",
        "menu_assessment_label": "🏢 Enterprise AI Readiness",
        "menu_assessment_sub": "Assess your org's AI readiness",
        "menu_compliance_label": "🛡️ AI Compliance & Traceability",
        "menu_compliance_sub": "Credibility Scoring & Compliance",
        "menu_dashboard_label": "📊 Strategic Dashboard",
        "menu_dashboard_sub": "Decision Priority Matrix",
        # Industry Tracks
        "track_general": "General Industry",
        "track_ai": "AI Market Research",
        "track_consulting": "Strategy Consulting",
        "track_risk": "Enterprise Risk Management",
        "track_manufacture": "Manufacturing EU Go-to-Market",
        "track_esg": "ESG / Sustainability",
        # Function 1: Report Search
        "search_input_tip": "Enter industry/track keywords, more details bring more accurate results",
        "search_btn": "Start Search",
        "search_loading": "Searching compliant industry reports...",
        "search_kw_empty": "Please enter search keywords",
        "search_pub_org": "🏢 Publisher: ",
        "search_pub_year": "📅 Publish Year: ",
        "search_abstract": "📄 Abstract: ",
        "search_credibility": "🔍 Credibility Score: ",
        # Function 2: Doc Summary
        "summary_mode": "Select Analysis Mode",
        "summary_mode_general": "General Document Summary",
        "summary_mode_indicator": "Industry Research Indicator Extraction",
        "summary_upload_tip": "Upload TXT/DOCX document, financial report, white paper",
        "summary_analyze_btn": "Start Analysis",
        "summary_analyze_loading": "Executing {mode}...",
        "summary_original_preview": "Original Document Preview",
        "summary_result_title": "✅ {mode} Result",
        "summary_download_btn": "📍 Download Word Result",
        "summary_download_filename": "Document_{mode}_Result.docx",
        # Function 3: Report Generation
        "generate_track_select": "Select Vertical Track Template",
        "generate_name_input": "Enter target industry/track/product name",
        "generate_ref_tip": "【Optional】Upload reference materials (generated content matches format first)",
        "generate_ref_upload": "Upload reference TXT/DOCX document",
        "generate_ref_preview": "Preview Reference Content",
        "generate_btn": "📝 Generate Consulting Report",
        "generate_loading": "Generating report for {track} track...",
        "generate_name_empty": "Please enter target industry/track name",
        "generate_ref_rule": "【Reference Rule】Generated content must prioritize the format from reference below:",
        "generate_report_title": "✉️ {name} | {track} Track Report",
        "generate_download_word": "📍 Download Word Report",
        "generate_download_ppt": "📊 Download Consulting PPT",
        "generate_word_filename": "{name}_{track}_Industry_Report.docx",
        "generate_ppt_filename": "{name}_{track}_Industry_Report.pptx",
        "generate_compliance_checkbox": "⚖️ Add Legal Compliance Risk Analysis Module",
        # Function 4: Multi-doc Compare
        "compare_tip": "Support 2-5 documents of the same track to generate strategic consulting comparative analysis report",
        "compare_upload_tip": "Upload TXT/DOCX documents for comparison",
        "compare_btn": "📈 Generate Comparative Analysis",
        "compare_loading": "Parsing documents and generating report...",
        "compare_file_min": "Please upload at least 2 documents for comparison",
        "compare_result_title": "✅ Comparative Analysis Report",
        "compare_download_word": "📍 Download Word Report",
        "compare_download_ppt": "📊 Download Consulting PPT",
        "compare_word_filename": "Track_Competitor_Analysis_Report.docx",
        "compare_ppt_filename": "Track_Competitor_Analysis_Report.pptx",
        "compare_priority_label": "📊 Include Strategic Priority Matrix",
        # Function 5: Template Rewrite
        "rewrite_flow": "Flow: Upload Template → Upload Target Document → One-click Rewrite → Preview → Download",
        "rewrite_template_upload": "1. Upload Template Document",
        "rewrite_content_upload": "2. Upload Target Document",
        "rewrite_template_preview": "Preview Template Content",
        "rewrite_content_preview": "Preview Target Content",
        "rewrite_btn": "✏️ Start Rewrite",
        "rewrite_loading": "Rewriting document with template style...",
        "rewrite_file_empty": "Please upload both template and target document first",
        "rewrite_result_title": "✅ Rewrite Result",
        "rewrite_download_word": "📍 Download Word Document",
        "rewrite_download_ppt": "📊 Download PPT Version",
        "rewrite_word_filename": "Document_Rewrite_Result.docx",
        "rewrite_ppt_filename": "Document_Rewrite_Result.pptx",
        # Function 6: Translation
        "translate_tip": "Support direct text translation or batch translation via TXT/DOCX upload",
        "translate_target_lang": "Select Target Language",
        "translate_mode": "Translation Mode",
        "translate_mode_text": "Direct Text Input",
        "translate_mode_file": "Upload Document",
        "translate_textarea_tip": "Enter business document content to translate",
        "translate_upload_tip": "Upload TXT/DOCX document to translate",
        "translate_original_preview": "Preview Original Content",
        "translate_btn": "🌐 Start Translation",
        "translate_loading": "Translating, please wait...",
        "translate_content_empty": "Please enter content or upload a valid document",
        "translate_result_title": "✅ Translation Result",
        "translate_download_btn": "📍 Download Word Result",
        "translate_download_filename": "Business_Document_Translation.docx",
        # Function 7: PDF to Word
        "pdf2word_tip": "Upload PDF → AI fix line breaks → Restore neat layout → Preview & Download",
        "pdf2word_upload_tip": "Upload PDF File",
        "pdf2word_loading": "Extracting PDF content and formatting with AI...",
        "pdf2word_preview_title": "📋 AI Formatted Content Preview",
        "pdf2word_download_word": "📍 Download Lossless Word",
        "pdf2word_download_ppt": "📊 Download PPT Version",
        "pdf2word_word_filename": "PDF_Conversion_Result.docx",
        "pdf2word_ppt_filename": "PDF_Conversion_Result.pptx",
        # New: About
        "about_title": "About This Product",
        "about_subtitle": "Product Methodology & Design Philosophy",
        "about_problem": "I. Problem Definition",
        "about_solution": "II. Solution",
        "about_architecture": "III. Technical Architecture",
        "about_value": "IV. Business Value",
        "about_principles": "V. Design Principles",
        "about_content": """### I. Problem Definition
Business decision-makers face information overload, fragmented data, and high costs for professional research reports. Traditional industry research relies on manual collection and organization, which is inefficient and prone to missing critical insights.

### II. Solution
Xiaobei AI Assistant leverages large language models and professional prompt engineering to provide a one-stop AI business decision support platform, covering report search, document analysis, report generation, and competitive comparison — drastically lowering the barrier to professional research.

### III. Technical Architecture
- **LLM Layer**: Volcano Engine Doubao Model (doubao-seed-2-0-pro) for high-quality semantic understanding and generation
- **Application Layer**: Streamlit frontend framework, Python backend
- **Document Processing**: python-docx / python-pptx / PyPDF2 / pdf2docx
- **Compliance Mechanisms**: Data source credibility scoring, AI content traceability, compliance check toggles

### IV. Business Value
- Reduces consulting-grade industry analysis costs by 90%+
- Compresses report generation time from days to minutes
- Covers 7 core business decision scenarios
- Supports bilingual (Chinese/English) professional output

### V. Design Principles
- **Responsible AI**: All generated content is traceable and verifiable
- **Professionalism**: Strict adherence to consulting industry report standards
- **Cross-disciplinary Integration**: Incorporating legal compliance, ESG, and organizational change perspectives
- **Product Methodology Driven**: Complete closed loop from problem definition to value delivery""",
        # New: Enterprise AI Readiness
        "assessment_title": "Enterprise AI Readiness Assessment",
        "assessment_subtitle": "Evaluate your organization's AI readiness and get customized recommendations",
        "assessment_company": "Company Name",
        "assessment_industry": "Industry",
        "assessment_size": "Employee Count",
        "assessment_revenue": "Annual Revenue (CNY 10K)",
        "assessment_existing_ai": "Current AI Adoption Level",
        "assessment_btn": "🔍 Start Assessment",
        "assessment_loading": "Assessing your organization's AI readiness...",
        "assessment_result_title": "📋 AI Readiness Assessment Report",
        "assessment_sizes": ["<100", "100-500", "500-2000", "2000+"],
        "assessment_ai_levels": ["Not using AI", "Experimenting (e.g. ChatGPT)", "Applied in some business units", "Established AI team/platform"],
        "assessment_company_empty": "Please enter company name",
        "assessment_fields": {
            "digital_level": "Digital Maturity Assessment",
            "ai_replace": "AI-Replaceable Task Inventory",
            "org_advice": "Organizational Restructuring Recommendations",
            "roi_estimate": "ROI Estimation"
        },
        # New: AI Compliance & Traceability (interactive)
        "compliance_title": "🛡️ AI Compliance & Traceability",
        "compliance_subtitle": "Responsible AI: Compliance Check + Traceability Report",
        "compliance_check_title": "I. Compliance Check (Paste Text for Review)",
        "compliance_check_tip": "Paste AI-generated or business document text below for compliance review",
        "compliance_check_placeholder": "Paste text to check here...",
        "compliance_check_btn": "🔍 Run Compliance Check",
        "compliance_check_loading": "Reviewing text for compliance issues...",
        "compliance_check_result_title": "✅ Compliance Check Result",
        "compliance_trace_title_new": "II. Traceability Report Generator",
        "compliance_trace_tip": "Enter a report topic, AI will generate a fully source-annotated report with references",
        "compliance_trace_placeholder": "e.g. Generative AI Industry Report 2024",
        "compliance_trace_btn": "📋 Generate Traceability Report",
        "compliance_trace_loading": "Generating traceability report with source annotations...",
        "compliance_trace_result_title": "✅ Traceability Report",
        "compliance_download_word": "📍 Download Word Report",
        "compliance_download_ppt": "📊 Download PPT Report",
        # Static compliance reference (placed below interactive section)
        "compliance_static_title": "III. Compliance Mechanism Reference",
        "compliance_score_title": "Data Source Credibility Scoring System",
        "compliance_score_desc": "Classify and score the credibility of data sources cited in AI-generated content to ensure reliability and verifiability.",
        "compliance_trace_title": "AI Content Traceability Chain",
        "compliance_trace_desc": "Establish a complete traceability chain from user input → prompt engineering → model inference → final output, with every step auditable.",
        "compliance_switch_title": "Compliance Check Toggles",
        "compliance_switch_desc": "Users can enable/disable compliance check items before generating reports, balancing efficiency and compliance needs.",
        "compliance_source_levels": "Credibility Tiers",
        "compliance_source_gov": "Government / Regulatory Bodies",
        "compliance_source_listed": "Listed Company Official Disclosure",
        "compliance_source_consulting": "Top Consulting / Research Institutions",
        "compliance_source_media": "Authoritative Financial / Tech Media",
        "compliance_source_unknown": "Unlabeled Source",
        "compliance_trace_flow": "User Input → Prompt Engineering (Compliance Rules Injected) → Doubao LLM Inference → Compliance Post-processing (Credibility Check) → Final Output",
        "compliance_switch_items": [
            "AI Content Compliance Check (auto-detect fabricated data, improper citations)",
            "Data Source Traceability (require AI to label source of each data point)",
            "Bias & Hallucination Detection (cross-validate key facts)"
        ],
        # New: Dashboard (interactive)
        "dashboard_title": "📊 Strategic Decision Dashboard",
        "dashboard_subtitle": "AI-generated Eisenhower Matrix & Priority Ranking",
        "dashboard_intro": "Enter your projects/initiatives below. AI will generate an Eisenhower Matrix (4-quadrant priority framework) and output P0-P3 execution priorities.",
        "dashboard_projects_tip": "Enter your projects / initiatives (one per line)",
        "dashboard_projects_placeholder": "e.g.\nQ3 Product Launch Campaign\nOfficial Website Redesign Project\nCompetitor Pricing Analysis Report\nCRM System Upgrade Evaluation",
        "dashboard_btn": "📊 Generate Decision Matrix",
        "dashboard_loading": "AI is generating your Eisenhower Matrix...",
        "dashboard_result_title": "✅ Strategic Decision Matrix",
        "dashboard_download_word": "📍 Download Word Report",
        "dashboard_download_ppt": "📊 Download PPT Report",
        # Others
        "func_not_found": "This function is not available yet",
        "ppt_title_default": "Consulting Report",
        "ppt_footer": "Xiaobei Consulting AI Assistant\nCompliant Generation | Traceable Data",
        "ppt_end_page": "End of Report",
        "ai_empty_response": "AI generated empty content, please try again~",
    }
}

# ==============================================================================
# 🌐 翻译目标语言选项
# ==============================================================================
TARGET_LANG_OPTIONS = {
    "zh": ["简体中文", "English", "日本語", "한국어", "繁体中文"],
    "en": ["Simplified Chinese", "English", "Japanese", "Korean", "Traditional Chinese"]
}

# ==============================================================================
# 🎯 用户核心配置区
# ==============================================================================
LOCAL_CONFIG = {
    "base_url": "https://ark.cn-beijing.volces.com/api/v3",
    "api_key": "ark-fc3c7e9f-d50d-48f5-8698-4955a37db662-5b27a",
    "model_name": "doubao-seed-2-0-pro-260215"
}

# ==============================================================================
# 📝 提示词配置区 — 提取所有硬编码Prompt到独立常量区
# ==============================================================================
PROMPT_CONFIG = {
    # ── 合规基础规则 ──
    "compliance_rule": """
    【强制合规要求，必须严格遵守】
    1. 所有数据、市场规模、增速、市场份额等量化内容，必须标注权威数据来源，包括但不限于：欧睿、IDC、乘联会、国家统计局、行业协会、上市公司财报、海关总署、贝恩/麦肯锡/波士顿咨询等权威机构发布的报告
    2. 绝对禁止虚构、编造任何数据、机构、事件、案例，所有内容必须符合行业真实情况
    3. 所有观点必须有对应的事实和数据支撑，禁止无依据的主观判断
    4. 严格遵循咨询行业报告的专业规范、结构逻辑和专业术语，语言正式、严谨、客观
    """,

    # ── 可信度评分规则（报告搜索增强） ──
    "credibility_scoring_rule": """
    【数据来源可信度评分规则】
    对每条报告的发布机构按以下标准标注可信度评分（满分10分）：
    - 政府/监管机构/国家统计局/央行/证监会：9-10分
    - 上市公司官方年报/招股书/公告：8-9分
    - 头部咨询机构（麦肯锡/BCG/贝恩/四大会计师事务所）：7-8分
    - 权威研究机构（IDC/Gartner/欧睿/艾瑞）：7-8分
    - 行业协会/券商研究所：6-7分
    - 财经媒体/科技媒体：4-6分
    - 未标注来源：1-3分
    """,

    # ── 行业报告模板 ──
    "formatting_requirement": """
【输出格式强制要求】
1. 必须使用Markdown格式输出，标题层级清晰：# 一级标题、## 二级标题、### 三级标题
2. 所有数据对比、指标汇总必须使用Markdown表格（| 列1 | 列2 |格式）
3. 所有列举项使用有序列表（1. 2. 3.）或无序列表（- 或 *）
4. 关键结论用**粗体**标注，重要数据用**粗体**突出
5. 报告结构完整：必须有目录（## 目录）、各章节、结论
6. 禁止输出纯文本堆积，必须分段、分点、分表格呈现
""",
    "industry_report_general": """
为【{name}】生成专业、合规的咨询级行业报告，必须严格遵守以下要求：
1. 报告结构必须包含7个核心部分：①行业定义与分类 ②市场规模与增长趋势 ③产业链上下游分析 ④竞争格局与核心玩家 ⑤用户画像与需求分析 ⑥行业痛点与发展趋势 ⑦投资机会与风险建议
2. {compliance_rule}
3. 结构清晰，段落分明，标题层级明确，符合正式咨询报告的排版规范
{formatting_requirement}
""",
    "industry_report_ai": """
为【{name}】生成AI领域专业市场研究报告，严格遵守AI行业研究规范，必须包含：
1. 核心结构：①赛道定义与技术路径 ②市场规模与投融资情况 ③技术成熟度与落地场景 ④核心厂商与竞争格局 ⑤政策监管环境 ⑥技术趋势与商业化痛点 ⑦市场机会与风险提示
2. {compliance_rule}
3. 重点突出AI技术落地的商业价值、市场竞争壁垒、客户付费意愿，符合一级市场AI赛道研究的专业规范
{formatting_requirement}
""",
    "industry_report_consulting": """
为【{name}】生成战略咨询级行业研究报告，严格遵循顶级咨询公司报告规范，必须包含：
1. 核心结构：①行业宏观环境（PEST分析）②市场规模与增长预测 ③产业链价值分布与利润池分析 ④竞争格局与五力模型分析 ⑤标杆企业战略与商业模式拆解 ⑥行业关键成功要素 ⑦企业进入战略与增长路径建议
2. {compliance_rule}
3. 重点突出战略洞察、可落地的商业建议，符合战略咨询项目的交付标准，逻辑严谨，洞察深刻
{formatting_requirement}
""",
    "industry_report_risk": """
为【{name}】生成企业风险管理视角的行业分析报告，严格遵循企业全面风险管理规范，必须包含：
1. 核心结构：①行业基本情况与经营环境 ②行业核心风险点识别（市场风险、信用风险、运营风险、合规风险、政策风险）③风险传导路径分析 ④行业标杆企业风险管理实践 ⑤风险应对策略与缓释措施 ⑥行业风险预警指标体系
2. {compliance_rule}
3. 重点突出风险的量化分析、发生概率与影响程度，符合企业内控与风险管理的专业要求
{formatting_requirement}
""",
    "industry_report_manufacture": """
为【{name}】生成制造业出海欧洲市场的专项分析报告，严格遵守跨境贸易与出海咨询规范，必须包含：
1. 核心结构：①欧洲目标市场准入政策与合规要求 ②市场规模与消费需求特征 ③欧洲本地竞争格局 ④供应链与物流方案分析 ⑤关税与税务筹划要点 ⑥本地化运营策略 ⑦出海风险与应对建议
2. {compliance_rule}
3. 重点突出欧洲市场合规要求、本地化运营难点、跨境供应链解决方案，符合制造业出海的真实业务需求
{formatting_requirement}
""",
    "industry_report_esg": """
为【{name}】生成ESG/可持续发展领域的专业分析报告，严格遵循GRI/TCFD/SASB等国际ESG披露标准，必须包含：
1. 核心结构：①ESG行业概述与政策环境 ②环境（E）维度分析——碳排放/能源结构/绿色技术 ③社会（S）维度分析——员工权益/供应链责任/社区影响 ④治理（G）维度分析——董事会结构/商业道德/信息披露 ⑤行业ESG评级对标（MSCI/标普/商道融绿）⑥ESG投资趋势与融资机会 ⑦可持续发展战略建议
2. {compliance_rule}
3. 重点突出量化ESG指标、国际评级对标、可落地的可持续发展路径，符合ESG投资与信息披露的专业要求
{formatting_requirement}
""",

    # ── 合规风险模块（附加到行业报告生成） ──
    "compliance_risk_module": """
    \n\n【额外模块：法律合规风险提示】\n请在本报告末尾增加独立的"法律合规风险提示"章节，包含以下内容：\n1. 行业关键法律法规与监管政策梳理（近3年）\n2. 数据安全与个人信息保护合规要求\n3. 反垄断与公平竞争合规要点\n4. 跨境经营法律风险（如涉及）\n5. 知识产权保护合规建议\n6. 劳动用工合规要点\n请确保每条合规风险提示都有对应的法律/法规依据来源。\n""",

    # ── 战略优先级矩阵指令（附加到对比分析） ──
    "priority_matrix_instruction": """
    \n\n【强制输出：战略建议优先级矩阵】\n请在对比分析报告的"战略建议"章节末尾，额外输出一个"战略建议优先级矩阵"，采用以下格式：\n\n| 优先级 | 战略建议 | 紧迫度(1-10) | 影响力(1-10) | 实施难度 | 建议时间窗口 |\n|--------|---------|-------------|-------------|---------|-------------|\n| P0-立即行动 | ... | 9-10 | 9-10 | 低-中 | 1个月内 |\n| P1-短期推进 | ... | 7-8 | 7-8 | 中 | 1-3个月 |\n| P2-中期规划 | ... | 5-6 | 5-6 | 中-高 | 3-6个月 |\n| P3-长期布局 | ... | 1-4 | 1-4 | 高 | 6个月以上 |\n\n请基于对比分析结果，在每个优先级下至少列出2条具体的、可执行的战略建议。\n""",

    # ── 文档总结 ──
    "doc_summary_general": """
对以下文档内容进行专业总结，核心输出4部分：1. 文档核心观点 2. 关键数据与信息 3. 行业竞争格局 4. 未来趋势与风险提示
{compliance_rule}
{formatting_requirement}
文档内容：{text}
""",
    "doc_summary_indicator": """
    对以下财报/行业白皮书/行研报告内容，进行行研核心指标提取，严格遵守以下要求：
    1. 必须提取的核心指标：市场规模、年复合增长率、市场集中度CR5/CR10、行业平均毛利率、核心竞品市场份额、核心财务指标、政策关键节点
    2. 所有指标必须标注对应的来源、统计年份、统计口径
    3. 最终输出必须是**标准Markdown结构化表格**，表格列名：指标名称、指标数值、统计周期、数据来源、备注说明
    4. 禁止虚构任何指标，无明确数据的指标标注「文档未提及」即可
    5. 表格输出完成后，补充100字以内的核心指标洞察总结
    文档内容：{text}
    """,

    # ── 报告搜索（增强：可信度评分） ──
    "report_search": """
    关键词：{keyword}，返回10条真实存在的行业报告，严格遵守格式要求：标题|发布机构|发布年份|可信度评分(1-10)|核心摘要
    {compliance_rule}
    {credibility_rule}
    禁止输出链接、网址、虚构内容，每条报告必须真实可查。可信度评分依据发布机构的权威性（政府/上市公司/咨询机构分别标注并加权）。
    """,

    # ── 模板改写 ──
    "template_rewrite": """
    你是专业咨询文档改写助手。请严格按以下流程执行：

    【第一步：分析模板文档】
    认真阅读【模板文档】，提取以下要素并严格模仿：
    - 文风特点：正式/半正式/口语化？论述型/清单型/叙事型？
    - 段落结构：每段长度、段落间逻辑关系（总-分、对比、递进等）
    - 标题风格：标题层级、是否编号、标题长度和措辞风格
    - 排版习惯：是否使用列表、表格、引用块？段落间距如何？
    - 开头结尾：模板如何开篇（背景引入/直接结论/提问式）？如何收尾（总结/展望/行动建议）？

    【第二步：按模板重写】
    把【待改写内容】按照第一步分析出的模板风格完整重写：
    1. 完全模仿模板的文风、结构、标题层级、段落格式、语气
    2. 不改变原文核心意思、核心数据、核心观点
    3. 优化语句的专业度、严谨性
    4. 改写后的文档应与模板"看起来像同一系列"
    5. {compliance_rule}

    【输出规则】
    - 不要添加任何多余解释、说明文字、装饰符号
    - 只输出改写后的完整内容
    - 改写结果必须是可直接使用的最终文档

    【模板文档】：
    {template_content}

    【待改写内容】：
    {original_content}
    """,

    # ── 文档翻译 ──
    "doc_translate": """
    你是专业商务文档翻译专家，严格遵守翻译要求：
    1. 目标语言：{target_lang}，严格按照目标语言进行专业翻译
    2. 精准翻译行业专业术语、商务表达、金融/咨询专业词汇，符合目标语言的正式商务文档规范
    3. 严格保留原文的段落结构、标题层级、表格格式，不改变原文核心意思、核心数据
    4. 翻译流畅自然，符合目标语言的商务写作习惯，无语法错误
    5. 不要添加任何额外解释、备注，只输出翻译后的完整内容

    需要翻译的原文：
    {text}
    """,

    # ── PDF排版 ──
    "pdf_format": """
    你是专业文档排版整理助手，请对下面PDF提取的乱序文字做无损规整排版：
    要求：
    1. 严格保留原文所有内容、核心数据、观点，不删减、不修改原文意思
    2. 按照原文的逻辑结构重新分段、换行、区分标题和正文，还原标题层级
    3. 修复PDF自动拆行、断句、乱换行、乱码问题
    4. 排版整洁、段落清晰、格式规范，适合直接保存为Word/PPT
    5. 不要加多余解释、序号、装饰符号，只输出规整后的完整内容

    需要整理的PDF原文：
    {text}
    """,

    # ── 多文档对比分析（增强：优先级矩阵） ──
    "multi_doc_compare": """
你是专业战略咨询顾问，基于以下上传的多份同赛道行业报告/竞品财报/白皮书，生成专业的对比分析报告，严格遵守要求：
1. 报告核心结构：①分析背景与对比范围 ②核心指标横向对比（市场规模、增速、盈利能力、市场份额等，输出结构化表格）③竞争格局与商业模式对比 ④核心优劣势差异分析 ⑤赛道机会与风险提示 ⑥战略建议{priority_instruction}
2. {compliance_rule}
3. 所有对比内容必须基于上传的文档内容，禁止添加文档外的虚构信息，重点突出核心差异与战略洞察
4. 结构清晰，符合战略咨询项目对比分析报告的专业规范
{formatting_requirement}

上传的文档内容合集：
{all_doc_text}
""",

    # ── 企业AI就绪度评估 ──
    "ai_readiness_assessment": """
你是一位资深的数字化转型与AI战略顾问。请基于以下企业信息，生成一份专业的"企业AI就绪度评估报告"。

【企业信息】
- 公司名称：{company_name}
- 所属行业：{industry}
- 员工规模：{company_size}
- 年营收：{revenue}万元人民币
- 现有AI应用情况：{ai_level}

请生成以下四个部分的评估内容：

## 一、数字化水平评估
基于公司规模、行业特征和现有AI应用情况，评估该企业的数字化成熟度（1-10分），并给出详细的优劣势分析。

## 二、AI可替代工作清单
列出该企业中可能被AI技术替代或增强的5-10项具体工作/流程，按替代优先级排序，每项标注：工作名称、当前人工耗时（估）、AI替代程度（%）、预期效率提升。

## 三、组织架构调整建议
为有效推进AI应用，建议组织架构需要做哪些调整（如建立AI卓越中心、设置首席AI官、调整部门职责等），每条建议标注优先级和预期影响。

## 四、ROI预估
基于行业平均水平，给出该企业AI应用的3年ROI预估值（节省的人力成本+效率提升带来的增收），分年度展示预估投入与收益。

{compliance_rule}
{formatting_requirement}
输出格式为结构化Markdown，语言正式专业，适合提交给董事会或管理层审阅。""",

    # ── 战略决策仪表盘：艾森豪威尔矩阵生成 ──
    "eisenhower_matrix": """
你是资深战略顾问。请基于以下用户输入的项目/方案列表，生成艾森豪威尔矩阵（Eisenhower Matrix）的战略优先级分析。

【用户输入的项目/方案】：
{projects_input}

【评分维度说明】：
- 紧迫度（1-10）：时间窗口的紧迫程度
- 影响力（1-10）：对战略目标的影响程度
- 实施难度（1-10）：资源/技术/组织阻力的综合难度

请按以下格式输出：

## 艾森豪威尔矩阵

### 🔴 第一象限：紧急且重要（P0 - 立即执行）
| 项目/方案 | 紧迫度 | 影响力 | 实施难度 | 建议行动 |
|-----------|--------|--------|---------|---------|
（列出所有紧迫度≥7且影响力≥7的项目）

### 🟡 第二象限：重要不紧急（P1 - 制定计划）
| 项目/方案 | 紧迫度 | 影响力 | 实施难度 | 建议行动 |
|-----------|--------|--------|---------|---------|
（列出所有紧迫度<7且影响力≥7的项目）

### 🟠 第三象限：紧急不重要（P2 - 委派/批量处理）
| 项目/方案 | 紧迫度 | 影响力 | 实施难度 | 建议行动 |
|-----------|--------|--------|---------|---------|
（列出所有紧迫度≥7且影响力<7的项目）

### 🟢 第四象限：不紧急不重要（P3 - 精简/延后）
| 项目/方案 | 紧迫度 | 影响力 | 实施难度 | 建议行动 |
|-----------|--------|--------|---------|---------|
（列出所有紧迫度<7且影响力<7的项目）

## 综合建议
基于矩阵结果，给出3-5条优先级排序的执行建议（按P0→P1→P2→P3顺序）。

{formatting_requirement}
""",

    # ── 合规检查：检查文本中的无来源断言 ──
    "compliance_check": """
你是AI内容合规审查专家。请对用户提供的文本进行"数据来源可信度"审查，严格遵守以下要求：

1. 逐条识别文本中**无明确数据来源**的断言、数据、结论
2. 对每条可疑内容标注：
   - 可疑内容原文（引用）
   - 问题类型：无来源数据 / 无来源结论 / 夸大表述 / 无法验证的断言
   - 风险等级：高（可能影响决策）/ 中（需要补充来源）/ 低（表述不够严谨）
   - 修改建议：如何补充来源或修正表述
3. 最后给出**整体合规评分**（0-100分）和**改进建议清单**

【待审查文本】：
{text_to_check}

{formatting_requirement}
""",

    # ── 溯源报告生成 ──
    "traceability_report": """
你是负责任AI内容生成专家。请基于以下报告主题，生成一份**带完整来源标注的溯源报告框架**。

【报告主题】：{report_topic}

请生成以下内容：

## 一、报告核心论点与所需数据来源清单
列出本报告需要引用的核心论点，以及每个论点对应的：
- 所需数据类型（市场规模/增速/竞争格局/政策文件等）
- 推荐权威来源（政府统计/上市公司/咨询机构/学术期刊等）
- 来源可信度评级（高/中/低）

## 二、溯源标注规范
规定本报告正文中如何标注来源（如：¹ 国家统计局2024年数据），并给出示例。

## 三、报告正文（带溯源标注）
按照行业报告规范，生成【报告主题】的完整报告正文，所有数据、结论后面必须用上标标注来源编号，文末附"参考文献"清单。

## 四、合规声明
声明本报告的数据来源、生成方式、局限性。

{compliance_rule}
{formatting_requirement}
""",
}

# ==============================================================================
# 📊 数据来源可信度评分映射表（常量区）
# ==============================================================================
CREDIBILITY_MAP = {
    "政府/监管机构": (9, 10),
    "国家统计局": (9, 10),
    "央行": (9, 10),
    "证监会": (9, 10),
    "上市公司年报": (8, 9),
    "上市公司公告": (8, 9),
    "麦肯锡": (8, 9),
    "波士顿咨询": (8, 9),
    "贝恩": (8, 9),
    "四大会计师事务所": (7, 8),
    "IDC": (7, 8),
    "Gartner": (7, 8),
    "欧睿": (7, 8),
    "艾瑞": (6, 7),
    "行业协会": (6, 7),
    "券商研究所": (6, 7),
    "财经媒体": (4, 6),
    "科技媒体": (4, 6),
}

# ==============================================================================
# 🌿 ESG/可持续发展设计原则（常量区）
# ==============================================================================
ESG_PRINCIPLES = """
ESG报告设计遵循以下国际标准框架：
- GRI（全球报告倡议组织）标准：通用披露 + 议题专项披露
- TCFD（气候相关财务信息披露工作组）：治理/战略/风险管理/指标与目标
- SASB（可持续会计准则委员会）：行业特定ESG指标
- 联合国SDGs（可持续发展目标）对标
"""

# ==============================================================================
# 📄 页面配置
# ==============================================================================
PAGE_CONFIG = {
    "page_icon": "😆"
}

# ==============================================================================
# 🚀 会话状态初始化
# ==============================================================================
def init_session_state():
    if "language" not in st.session_state:
        st.session_state.language = "zh"
    if "selected_tab" not in st.session_state:
        st.session_state.selected_tab = ""
    if "rewrite_result" not in st.session_state:
        st.session_state.rewrite_result = ""
    if "rewrite_generating" not in st.session_state:
        st.session_state.rewrite_generating = False
    if "translate_result" not in st.session_state:
        st.session_state.translate_result = ""
    if "translate_generating" not in st.session_state:
        st.session_state.translate_generating = False
    if "compare_result" not in st.session_state:
        st.session_state.compare_result = ""
    if "compare_generating" not in st.session_state:
        st.session_state.compare_generating = False
    if "assessment_result" not in st.session_state:
        st.session_state.assessment_result = ""

# 初始化会话状态
init_session_state()

# 获取当前语言设置
if "language" not in st.session_state:
    st.session_state.language = "zh"
current_lang = st.session_state.language
lang = LANG_PACK[current_lang]

# ==============================================================================
# 📋 菜单配置 — 7原有 + 4新增
# ==============================================================================
MENU_CONFIG = [
    {"id": "search",        "label": lang["menu_search_label"],        "sub_title": lang["menu_search_sub"]},
    {"id": "summary",       "label": lang["menu_summary_label"],       "sub_title": lang["menu_summary_sub"]},
    {"id": "generate",      "label": lang["menu_generate_label"],      "sub_title": lang["menu_generate_sub"]},
    {"id": "compare",       "label": lang["menu_compare_label"],       "sub_title": lang["menu_compare_sub"]},
    {"id": "rewrite",       "label": lang["menu_rewrite_label"],       "sub_title": lang["menu_rewrite_sub"]},
    {"id": "translate",     "label": lang["menu_translate_label"],     "sub_title": lang["menu_translate_sub"]},
    {"id": "pdf2word",      "label": lang["menu_pdf2word_label"],      "sub_title": lang["menu_pdf2word_sub"]},
    {"id": "dashboard",     "label": lang["menu_dashboard_label"],     "sub_title": lang["menu_dashboard_sub"]},
    {"id": "assessment",    "label": lang["menu_assessment_label"],    "sub_title": lang["menu_assessment_sub"]},
    {"id": "compliance",    "label": lang["menu_compliance_label"],    "sub_title": lang["menu_compliance_sub"]},
    {"id": "about",         "label": lang["menu_about_label"],         "sub_title": lang["menu_about_sub"]},
]

# ==============================================================================
# 🏭 赛道配置 — 5原有 + ESG新增
# ==============================================================================
INDUSTRY_TRACKS = [
    lang["track_general"],
    lang["track_ai"],
    lang["track_consulting"],
    lang["track_risk"],
    lang["track_manufacture"],
    lang["track_esg"],
]

TRACK_PROMPT_MAP = {
    lang["track_general"]:      PROMPT_CONFIG["industry_report_general"],
    lang["track_ai"]:           PROMPT_CONFIG["industry_report_ai"],
    lang["track_consulting"]:   PROMPT_CONFIG["industry_report_consulting"],
    lang["track_risk"]:         PROMPT_CONFIG["industry_report_risk"],
    lang["track_manufacture"]:  PROMPT_CONFIG["industry_report_manufacture"],
    lang["track_esg"]:          PROMPT_CONFIG["industry_report_esg"],
}

MENU_LABELS = [item["label"] for item in MENU_CONFIG]
MENU_MAP = {item["label"]: item for item in MENU_CONFIG}

client = OpenAI(
    base_url=os.getenv("ARK_BASE_URL", LOCAL_CONFIG["base_url"]),
    api_key=os.getenv("ARK_API_KEY", LOCAL_CONFIG["api_key"]),
)

# ==============================================================================
# 🔧 通用工具函数
# ==============================================================================
def ai_request(prompt, target_lang=None):
    """调用火山引擎豆包API，支持目标语言控制"""
    try:
        if target_lang:
            if target_lang == "en":
                prompt = f"Please respond in English only. {prompt}"
            elif target_lang == "zh":
                prompt = f"请用中文回答。{prompt}"

        response = client.chat.completions.create(
            model=LOCAL_CONFIG["model_name"],
            messages=[{"role": "user", "content": prompt}]
        )
        full_text = response.choices[0].message.content
        return full_text.strip() if full_text else lang.get("ai_empty_response", "AI生成内容为空，请重试~")
    except Exception as e:
        st.error(f"{lang['error']}: {str(e)}")
        return ""

def read_file(uploaded_file):
    """读取上传的TXT/DOCX文件内容"""
    try:
        raw = uploaded_file.read()
        if uploaded_file.name.lower().endswith(".docx"):
            doc = Document(io.BytesIO(raw))
            full_text = "\n".join([p.text for p in doc.paragraphs])
            return full_text.strip()
        elif uploaded_file.name.lower().endswith(".txt"):
            return raw.decode("utf-8", errors="ignore").strip()
        else:
            return "不支持的文件格式"
    except Exception as e:
        return f"文件读取失败：{str(e)}"

def clean_markdown(text):
    """清理AI输出中的Markdown格式符号，保留纯文本内容"""
    import re
    # 移除粗体/斜体标记 **text** 和 *text*
    text = re.sub(r'\*{1,3}([^\*\n]+)\*{1,3}', r'\1', text)
    # 移除行内代码标记 `text`
    text = re.sub(r'`([^`\n]+)`', r'\1', text)
    # 移除标题标记 # （保留标题文字）
    text = re.sub(r'^#{1,6}\s*', '', text, flags=re.MULTILINE)
    # 移除列表标记 - 或 * 或数字. （保留列表文字）
    text = re.sub(r'^[\s]*[-*+]\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'^[\s]*\d+\.\s+', '', text, flags=re.MULTILINE)
    # 移除引用标记 >
    text = re.sub(r'^>\s*', '', text, flags=re.MULTILINE)
    # 移除分隔线 --- 或 ***
    text = re.sub(r'^[-*_]{3,}\s*$', '', text, flags=re.MULTILINE)
    # 移除链接标记 [text](url) → text
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
    return text.strip()


def generate_word_file(content):
    """将Markdown内容转换为富格式Word文档（保留标题层级、表格、列表、粗体斜体）"""
    import re
    from docx.oxml import OxmlElement
    from docx.shared import Pt as DocxPt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # 设置默认样式
    style = doc.styles['Normal']
    style.font.name = '微软雅黑'
    style.font.size = DocxPt(11)
    style.paragraph_format.space_after = DocxPt(4)

    def _add_inline_runs(paragraph, text):
        """解析行内格式(**bold**, *italic*, `code`)并添加runs"""
        paragraph.clear()
        pattern = r'(\*\*\*(.+?)\*\*\*|\*\*(.+?)\*\*|\*(.+?)\*|`(.+?)`)'
        parts = re.split(pattern, text)
        for part in parts:
            if part is None or part == '':
                continue
            # 检查是否被粗体+斜体匹配(***...***)
            bold_italic_match = re.match(r'\*\*\*(.+?)\*\*\*$', part)
            if bold_italic_match is None:
                bold_italic_match = re.match(r'^\*\*\*(.+?)\*\*\*', part)
            if bold_italic_match:
                p = paragraph if paragraph.runs else paragraph
                continue  # handled by the triple-asterisk case below

            run = paragraph.add_run(part)
            run.font.name = '微软雅黑'
            run._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')

        # Second pass: find bold+italic triple patterns
        full_text = paragraph.text
        if '***' in full_text:
            paragraph.clear()
            bold_italic_parts = re.split(r'(\*\*\*.+?\*\*\*)', full_text)
            for segment in bold_italic_parts:
                if segment.startswith('***') and segment.endswith('***'):
                    run = paragraph.add_run(segment[3:-3])
                    run.bold = True
                    run.italic = True
                    run.font.name = '微软雅黑'
                    run._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
                elif '**' in segment:
                    sub_parts = re.split(r'(\*\*.+?\*\*)', segment)
                    for sp in sub_parts:
                        if sp.startswith('**') and sp.endswith('**'):
                            run = paragraph.add_run(sp[2:-2])
                            run.bold = True
                            run.font.name = '微软雅黑'
                            run._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
                        elif '*' in sp:
                            sub2 = re.split(r'(\*.+?\*)', sp)
                            for s2 in sub2:
                                if s2.startswith('*') and s2.endswith('*') and not s2.startswith('**'):
                                    run = paragraph.add_run(s2[1:-1])
                                    run.italic = True
                                    run.font.name = '微软雅黑'
                                    run._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
                                elif s2:
                                    run = paragraph.add_run(s2)
                                    run.font.name = '微软雅黑'
                                    run._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
                        else:
                            run = paragraph.add_run(sp)
                            run.font.name = '微软雅黑'
                            run._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
                else:
                    run = paragraph.add_run(segment)
                    run.font.name = '微软雅黑'
                    run._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')

    lines = content.split('\n')
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if not stripped:
            i += 1
            continue

        # ── 检测 Markdown 表格 ──
        if '|' in stripped and i + 1 < len(lines):
            next_line = lines[i + 1].strip()
            # 分隔行: 包含 | 、-、: 且不含其他字符
            if re.match(r'^[\|\s\-:]+$', next_line) and '|' in next_line:
                # 收集表头
                header_cells = [c.strip() for c in stripped.split('|')]
                header_cells = [c for c in header_cells if c]  # 去空首尾
                table_data = [header_cells]
                i += 2  # 跳过分隔行
                # 收集数据行
                while i < len(lines) and '|' in lines[i].strip() and lines[i].strip():
                    row_cells = [c.strip() for c in lines[i].strip().split('|')]
                    row_cells = [c for c in row_cells if c]
                    if row_cells:
                        table_data.append(row_cells)
                    i += 1
                # 创建 Word 表格
                if table_data:
                    max_cols = max(len(r) for r in table_data)
                    # 补齐列数
                    for r in table_data:
                        while len(r) < max_cols:
                            r.append('')
                    table = doc.add_table(rows=len(table_data), cols=max_cols, style='Light Grid Accent 1')
                    for r_idx, row_data in enumerate(table_data):
                        for c_idx, cell_text in enumerate(row_data):
                            if c_idx < max_cols:
                                cell = table.cell(r_idx, c_idx)
                                cell.text = cell_text
                                for p in cell.paragraphs:
                                    p.paragraph_format.space_before = DocxPt(2)
                                    p.paragraph_format.space_after = DocxPt(2)
                                    for run in p.runs:
                                        run.font.name = '微软雅黑'
                                        run.font.size = DocxPt(10)
                                        run._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
                    doc.add_paragraph()  # 表后间距
                continue

        # ── 检测标题 # ~ ###### ──
        heading_match = re.match(r'^(#{1,6})\s+(.+)$', stripped)
        if heading_match:
            level = min(len(heading_match.group(1)), 3)  # 最多 Heading 3
            heading_text = re.sub(r'\*{1,3}([^*\n]+)\*{1,3}', r'\1', heading_match.group(2))
            heading_text = re.sub(r'`([^`\n]+)`', r'\1', heading_text)
            h = doc.add_heading(heading_text, level=level)
            for run in h.runs:
                run.font.name = '微软雅黑'
                run._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
            i += 1
            continue

        # ── 检测无序列表 - 或 * ──
        bullet_match = re.match(r'^[\s]*[-*+]\s+(.+)$', stripped)
        if bullet_match:
            text = bullet_match.group(1)
            p = doc.add_paragraph(style='List Bullet')
            _add_inline_runs(p, text)
            i += 1
            continue

        # ── 检测有序列表 1. 或 1) ──
        num_match = re.match(r'^[\s]*(\d+)[.)]\s+(.+)$', stripped)
        if num_match:
            text = num_match.group(2)
            p = doc.add_paragraph(style='List Number')
            _add_inline_runs(p, text)
            i += 1
            continue

        # ── 普通段落 ──
        p = doc.add_paragraph()
        _add_inline_runs(p, stripped)
        i += 1

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def _set_chinese_font(shape):
    """递归设置Shape中所有文本的中文字体（微软雅黑）"""
    from pptx.util import Pt as _Pt
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.name = "微软雅黑"
            for run in paragraph.runs:
                run.font.name = "微软雅黑"
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = "微软雅黑"
                    for run in paragraph.runs:
                        run.font.name = "微软雅黑"


def _add_content_slide(prs, slide_layout, title_text, body_lines):
    """添加一页内容幻灯片，自动处理标题和正文，正文自动换页"""
    import re
    # 确定标题是否为标题样式
    is_real_title = bool(re.search(r'^第[一二三四五六七八九十\d]+[章节部分]|^[一二三四五六七八九十\d]+[、.]|'
                                   r'^[A-Ga-g]\d*[.、]|^[IVX]+[.、]|^引言|^前言|^概述|^总结|^结论|^附录|'
                                   r'^Abstract|^Introduction|^Conclusion|^Summary', title_text, re.UNICODE))

    lines_per_slide = 15  # 每页最多行数
    chunk_count = 0

    for chunk_start in range(0, len(body_lines), lines_per_slide):
        chunk = body_lines[chunk_start:chunk_start + lines_per_slide]
        if not chunk:
            continue
        slide = prs.slides.add_slide(slide_layout)
        # 设置标题
        title_shape = slide.shapes.title
        if chunk_count == 0:
            title_shape.text = title_text
        else:
            title_shape.text = f"{title_text}（续）"
        _set_chinese_font(title_shape)
        title_shape.text_frame.paragraphs[0].font.size = Pt(22)

        # 设置正文
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, line in enumerate(chunk):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(14)
                p.font.name = "微软雅黑"
                p.space_after = Pt(6)
        chunk_count += 1


def generate_ppt_file(content, title, footer, end_text):
    """生成专业演示文稿：封面→目录→章节分隔页→内容页→总结→感谢页，带页码和主题色"""
    import re
    from datetime import datetime
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN as PPT_ALIGN

    try:
        prs = Presentation()
        # 幻灯片尺寸设为宽屏 16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        THEME_COLOR = RGBColor(0x1A, 0x3C, 0x6E)  # 深蓝色主题
        ACCENT_COLOR = RGBColor(0xE8, 0x6A, 0x17)  # 橙色强调
        LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)

        slide_number = [0]  # 用列表实现闭包引用

        def _add_slide_number(slide):
            """在幻灯片右下角添加页码"""
            slide_number[0] += 1
            left = prs.slide_width - Inches(1.2)
            top = prs.slide_height - Inches(0.5)
            txBox = slide.shapes.add_textbox(left, top, Inches(1.0), Inches(0.4))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = str(slide_number[0])
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
            p.alignment = PPT_ALIGN.RIGHT

        def _add_section_divider(prs, section_title, section_number):
            """添加章节分隔页（深色背景 + 大标题）"""
            blank_layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(blank_layout)
            # 背景色矩形
            from pptx.util import Emu
            shape = slide.shapes.add_shape(
                1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = THEME_COLOR
            shape.line.fill.background()

            # 章节编号
            txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(1.0))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"PART {section_number}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0xAA, 0xBB, 0xDD)
            p.font.name = '微软雅黑'

            # 章节标题
            txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10), Inches(2.0))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = section_title
            p2.font.size = Pt(40)
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p2.font.name = '微软雅黑'

            _add_slide_number(slide)
            return slide

        # ═══════════════════════════════════════
        # Slide 1: 封面
        # ═══════════════════════════════════════
        blank_layout = prs.slide_layouts[6]
        cover = prs.slides.add_slide(blank_layout)
        # 深色背景条
        shape = cover.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(3.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = THEME_COLOR
        shape.line.fill.background()

        # 主标题
        txBox = cover.shapes.add_textbox(Inches(1.2), Inches(0.8), Inches(10.5), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.name = '微软雅黑'

        # 副标题
        txBox2 = cover.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.5), Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = footer.replace('\n', '  |  ')
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
        p2.font.name = '微软雅黑'

        # 日期
        txBox3 = cover.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(5), Inches(0.5))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = datetime.now().strftime("%Y年%m月%d日")
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p3.font.name = '微软雅黑'

        _add_slide_number(cover)

        # ═══════════════════════════════════════
        # 解析内容结构：找出章节标题（## 开头）
        # ═══════════════════════════════════════
        lines = content.split('\n')
        sections = []  # [(标题, 内容行列表), ...]
        current_section_title = None
        current_section_lines = []

        for line in lines:
            stripped = line.strip()
            h2_match = re.match(r'^##\s+(.+)$', stripped)
            h1_match = re.match(r'^#\s+(.+)$', stripped)
            if h1_match or h2_match:
                if current_section_title is not None or current_section_lines:
                    sections.append((current_section_title, current_section_lines))
                current_section_title = (h1_match or h2_match).group(1)
                current_section_title = re.sub(r'\*{1,3}([^*\n]+)\*{1,3}', r'\1', current_section_title)
                current_section_lines = []
            else:
                if stripped:
                    current_section_lines.append(stripped)

        # 保存最后一节
        if current_section_title is not None or current_section_lines:
            sections.append((current_section_title, current_section_lines))

        # 如果没有检测到章节标题，把全部内容作为一节
        if not sections or (len(sections) == 1 and sections[0][0] is None):
            sections = [(title, [l.strip() for l in lines if l.strip()])]

        # ═══════════════════════════════════════
        # Slide 2: 目录页
        # ═══════════════════════════════════════
        toc_slide = prs.slides.add_slide(blank_layout)
        txBox = toc_slide.shapes.add_textbox(Inches(1.5), Inches(0.8), Inches(10), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "目  录"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = THEME_COLOR
        p.font.name = '微软雅黑'

        # 目录项
        txBox2 = toc_slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(4.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for idx, (s_title, _) in enumerate(sections):
            if s_title:
                if idx == 0:
                    p = tf2.paragraphs[0]
                else:
                    p = tf2.add_paragraph()
                p.text = f"{idx + 1}.  {s_title}"
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                p.font.name = '微软雅黑'
                p.space_after = Pt(14)
        _add_slide_number(toc_slide)

        # ═══════════════════════════════════════
        # 各章节内容
        # ═══════════════════════════════════════
        title_content_layout = prs.slide_layouts[1]  # Title and Content

        for sec_idx, (sec_title, sec_lines) in enumerate(sections):
            if not sec_title:
                sec_title = title

            # 章节分隔页
            _add_section_divider(prs, sec_title, sec_idx + 1)

            # 清理内容：去除 markdown 格式符号，提取要点
            cleaned_lines = []
            for line in sec_lines:
                cl = re.sub(r'\*{1,3}([^*\n]+)\*{1,3}', r'\1', line)
                cl = re.sub(r'`([^`\n]+)`', r'\1', cl)
                cl = re.sub(r'^#{1,6}\s+', '', cl)
                cl = re.sub(r'^[\s]*[-*+]\s+', '', cl)
                cl = re.sub(r'^[\s]*\d+[.)]\s+', '', cl)
                cl = cl.strip()
                if cl:
                    cleaned_lines.append(cl)

            # 如果内容完全为空，跳过
            if not cleaned_lines:
                continue

            # 每页最多 7 个要点
            points_per_page = 7
            for page_start in range(0, len(cleaned_lines), points_per_page):
                page_lines = cleaned_lines[page_start:page_start + points_per_page]
                slide = prs.slides.add_slide(title_content_layout)

                # 标题
                ts = slide.shapes.title
                ts.text = sec_title if page_start == 0 else f"{sec_title}（续）"
                _set_chinese_font(ts)
                for run in ts.text_frame.paragraphs[0].runs:
                    run.font.size = Pt(24)
                    run.font.bold = True
                    run.font.color.rgb = THEME_COLOR

                # 内容区域
                if len(slide.placeholders) > 1:
                    tf = slide.placeholders[1].text_frame
                    tf.clear()
                    tf.word_wrap = True
                    for pi, line in enumerate(page_lines):
                        if pi == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = line
                        p.font.size = Pt(16)
                        p.font.name = '微软雅黑'
                        p.space_after = Pt(10)
                        p.level = 0
                        # 添加项目符号
                        p.level = 0

                _add_slide_number(slide)

        # ═══════════════════════════════════════
        # 总结页
        # ═══════════════════════════════════════
        summary_slide = prs.slides.add_slide(blank_layout)
        txBox = summary_slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(2.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "总结与展望"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = THEME_COLOR
        p.font.name = '微软雅黑'
        p.alignment = PPT_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = "感谢您的关注与支持"
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p2.font.name = '微软雅黑'
        p2.alignment = PPT_ALIGN.CENTER
        p2.space_before = Pt(20)

        _add_slide_number(summary_slide)

        # ═══════════════════════════════════════
        # 感谢页
        # ═══════════════════════════════════════
        thanks_slide = prs.slides.add_slide(blank_layout)
        shape = thanks_slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = THEME_COLOR
        shape.line.fill.background()

        txBox = thanks_slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(2.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = end_text
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.name = '微软雅黑'
        p.alignment = PPT_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = footer.replace('\n', '  |  ')
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0xAA, 0xBB, 0xDD)
        p2.font.name = '微软雅黑'
        p2.alignment = PPT_ALIGN.CENTER
        p2.space_before = Pt(16)

        _add_slide_number(thanks_slide)

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer

    except Exception as e:
        st.error(f"PPT生成失败: {str(e)}")
        # 兜底：生成最简PPT
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer

def extract_pdf_text(pdf_file):
    """从PDF文件中提取文本"""
    try:
        reader = PdfReader(pdf_file)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n\n"
        return text.strip() if text else "PDF文本提取为空"
    except Exception as e:
        return f"PDF读取失败：{str(e)}"

# ==============================================================================
# 📌 功能1：全网报告搜索（增强：数据来源可信度评分）
# ==============================================================================
def render_search():
    kw = st.text_input(lang["search_input_tip"])
    if st.button(lang["search_btn"], use_container_width=True):
        if kw:
            with st.spinner(lang["search_loading"]):
                prompt = PROMPT_CONFIG["report_search"].format(
                    keyword=kw,
                    compliance_rule=PROMPT_CONFIG["compliance_rule"],
                    credibility_rule=PROMPT_CONFIG["credibility_scoring_rule"]
                )
                content = ai_request(prompt, target_lang=st.session_state.language)
                if content:
                    lines = content.strip().split("\n")
                    for i, line in enumerate(lines):
                        if "|" in line:
                            p = line.split("|")
                            if len(p) >= 4:
                                title = p[0].strip()
                                org = p[1].strip()
                                year = p[2].strip()
                                credibility = p[3].strip() if len(p) >= 5 else ""
                                abstract = p[3].strip() if len(p) == 4 else (p[4].strip() if len(p) >= 5 else "")

                                # 构造真实搜索链接
                                search_query = f"{title} {org} {year} 报告 PDF"
                                google_url = f"https://www.google.com/search?q={search_query.replace(' ', '+')}"
                                baidu_url = f"https://www.baidu.com/s?wd={search_query.replace(' ', '%20')}"

                                st.markdown(f"### {i+1}. {title}")
                                st.write(f"{lang['search_pub_org']}{org} | {lang['search_pub_year']}{year}")

                                # 展示可信度评分
                                if credibility:
                                    try:
                                        score_val = float(credibility)
                                        if score_val >= 8:
                                            score_color = "green"
                                            score_label = "高可信度"
                                        elif score_val >= 5:
                                            score_color = "orange"
                                            score_label = "中等可信度"
                                        else:
                                            score_color = "red"
                                            score_label = "低可信度"
                                        st.markdown(
                                            f"{lang['search_credibility']} "
                                            f":{score_color}[**{credibility}/10** — {score_label}]"
                                        )
                                    except ValueError:
                                        st.markdown(f"{lang['search_credibility']}{credibility}")

                                st.write(f"{lang['search_abstract']}{abstract}")
                                st.markdown(f"🔗 [Google搜索]({google_url})  |  [百度搜索]({baidu_url})")
                                st.divider()
        else:
            st.warning(f"{lang['warning']}: {lang['search_kw_empty']}")

# ==============================================================================
# 📌 功能2：文档总结/数据提取（保持原样）
# ==============================================================================
def render_summary():
    summary_mode = st.radio(
        lang["summary_mode"],
        options=[lang["summary_mode_general"], lang["summary_mode_indicator"]],
        horizontal=True
    )
    st.markdown("---")
    f = st.file_uploader(lang["summary_upload_tip"], type=["txt","docx"])
    if f and st.button(lang["summary_analyze_btn"], use_container_width=True):
        with st.spinner(lang["summary_analyze_loading"].format(mode=summary_mode)):
            txt = read_file(f)
            st.text_area(lang["summary_original_preview"], txt, height=200)
            st.markdown("---")
            if summary_mode == lang["summary_mode_general"]:
                prompt = PROMPT_CONFIG["doc_summary_general"].format(text=txt[:3500], compliance_rule=PROMPT_CONFIG["compliance_rule"], formatting_requirement=PROMPT_CONFIG["formatting_requirement"])
            else:
                prompt = PROMPT_CONFIG["doc_summary_indicator"].format(text=txt[:6000], compliance_rule=PROMPT_CONFIG["compliance_rule"])
            res = ai_request(prompt, target_lang=st.session_state.language)
            st.markdown(f"### {lang['summary_result_title'].format(mode=summary_mode)}")
            st.write(res)
            word_buf = generate_word_file(res)
            st.download_button(
                label=lang["summary_download_btn"],
                data=word_buf,
                file_name=lang["summary_download_filename"].format(mode=summary_mode),
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )

# ==============================================================================
# 📌 功能3：行业报告生成（增强：ESG赛道 + 合规风险模块）
# ==============================================================================
def render_generate():
    col1, col2 = st.columns(2)
    with col1:
        selected_track = st.selectbox(lang["generate_track_select"], options=INDUSTRY_TRACKS)
    with col2:
        industry_name = st.text_input(lang["generate_name_input"])
    st.markdown("---")
    st.markdown(lang["generate_ref_tip"])
    reference_file = st.file_uploader(lang["generate_ref_upload"], type=["txt","docx"], key="reference_file")
    reference_text = ""
    if reference_file:
        reference_text = read_file(reference_file)
        with st.expander(lang["generate_ref_preview"]):
            st.text_area("参考资料", reference_text, height=200)

    # 新增：法律合规风险分析复选框
    st.markdown("---")
    add_compliance = st.checkbox(lang["generate_compliance_checkbox"], value=False)

    st.markdown("---")
    if st.button(lang["generate_btn"], use_container_width=True):
        if not industry_name:
            st.warning(f"{lang['warning']}: {lang['generate_name_empty']}")
        else:
            with st.spinner(lang["generate_loading"].format(track=selected_track)):
                base_prompt = TRACK_PROMPT_MAP[selected_track]
                full_prompt = base_prompt.format(
                    name=industry_name,
                    compliance_rule=PROMPT_CONFIG["compliance_rule"],
                    formatting_requirement=PROMPT_CONFIG["formatting_requirement"]
                )
                if reference_text:
                    full_prompt += f"\n\n{lang['generate_ref_rule']}\n{reference_text[:3000]}"
                # 附加合规风险模块
                if add_compliance:
                    full_prompt += PROMPT_CONFIG["compliance_risk_module"]

                report_content = ai_request(full_prompt, target_lang=st.session_state.language)
                st.markdown(f"### {lang['generate_report_title'].format(name=industry_name, track=selected_track)}")
                st.write(report_content)
                col_word, col_ppt = st.columns(2)
                with col_word:
                    word_buf = generate_word_file(report_content)
                    st.download_button(
                        label=lang["generate_download_word"],
                        data=word_buf,
                        file_name=lang["generate_word_filename"].format(name=industry_name, track=selected_track),
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        use_container_width=True
                    )
                with col_ppt:
                    ppt_buf = generate_ppt_file(
                        content=report_content,
                        title=f"{industry_name} {lang['ppt_title_default']}",
                        footer=lang["ppt_footer"],
                        end_text=lang["ppt_end_page"]
                    )
                    st.download_button(
                        label=lang["generate_download_ppt"],
                        data=ppt_buf,
                        file_name=lang["generate_ppt_filename"].format(name=industry_name, track=selected_track),
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )

# ==============================================================================
# 📌 功能4：多文档竞品对比分析（增强：战略优先级矩阵）
# ==============================================================================
def render_compare():
    st.markdown(lang["compare_tip"])

    # 新增：优先级矩阵开关
    include_priority = st.checkbox(lang["compare_priority_label"], value=True)

    upload_files = st.file_uploader(
        lang["compare_upload_tip"],
        type=["txt","docx"],
        accept_multiple_files=True
    )
    st.markdown("---")
    if st.button(lang["compare_btn"], use_container_width=True):
        if not upload_files or len(upload_files) < 2:
            st.warning(f"{lang['warning']}: {lang['compare_file_min']}")
        else:
            with st.spinner(lang["compare_loading"]):
                all_doc_text = ""
                for i, file in enumerate(upload_files):
                    file_text = read_file(file)
                    all_doc_text += f"===== 文档{i+1}：{file.name} =====\n{file_text[:3000]}\n\n"

                # 根据开关决定是否附加优先级矩阵指令
                priority_instruction = PROMPT_CONFIG["priority_matrix_instruction"] if include_priority else ""

                prompt = PROMPT_CONFIG["multi_doc_compare"].format(
                    all_doc_text=all_doc_text,
                    compliance_rule=PROMPT_CONFIG["compliance_rule"],
                    priority_instruction=priority_instruction,
                    formatting_requirement=PROMPT_CONFIG["formatting_requirement"]
                )
                compare_result = ai_request(prompt, target_lang=st.session_state.language)
                st.session_state.compare_result = compare_result

    if st.session_state.compare_result:
        st.markdown(f"### {lang['compare_result_title']}")
        st.write(st.session_state.compare_result)
        col_word, col_ppt = st.columns(2)
        with col_word:
            word_buf = generate_word_file(st.session_state.compare_result)
            st.download_button(
                label=lang["compare_download_word"],
                data=word_buf,
                file_name=lang["compare_word_filename"],
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )
        with col_ppt:
            ppt_buf = generate_ppt_file(
                content=st.session_state.compare_result,
                title=lang["compare_result_title"],
                footer=lang["ppt_footer"],
                end_text=lang["ppt_end_page"]
            )
            st.download_button(
                label=lang["compare_download_ppt"],
                data=ppt_buf,
                file_name=lang["compare_ppt_filename"],
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )

# ==============================================================================
# 📌 功能5：仿照模板改写文档（保持原样）
# ==============================================================================
def render_rewrite():
    st.markdown(lang["rewrite_flow"])
    col1, col2 = st.columns(2)
    template_text = ""
    content_text = ""
    with col1:
        template_file = st.file_uploader(lang["rewrite_template_upload"], type=["txt","docx"], key="template_file")
        if template_file:
            template_text = read_file(template_file)
            with st.expander(lang["rewrite_template_preview"]):
                st.text_area("Template", template_text, height=280, key="template_preview")
    with col2:
        content_file = st.file_uploader(lang["rewrite_content_upload"], type=["txt","docx"], key="content_file")
        if content_file:
            content_text = read_file(content_file)
            with st.expander(lang["rewrite_content_preview"]):
                st.text_area("Original", content_text, height=280, key="content_preview")
    st.markdown("---")
    if st.button(lang["rewrite_btn"], use_container_width=True, disabled=st.session_state.rewrite_generating):
        if not template_file or not content_file:
            st.warning(f"{lang['warning']}: {lang['rewrite_file_empty']}")
        else:
            st.session_state.rewrite_result = ""
            st.session_state.rewrite_generating = True
    if st.session_state.rewrite_generating and not st.session_state.rewrite_result:
        with st.spinner(lang["rewrite_loading"]):
            prompt = PROMPT_CONFIG["template_rewrite"].format(
                template_content=template_text[:4000],
                original_content=content_text[:5000],
                compliance_rule=PROMPT_CONFIG["compliance_rule"]
            )
            result_text = ai_request(prompt, target_lang=st.session_state.language)
            st.session_state.rewrite_result = result_text
            st.session_state.rewrite_generating = False
    if st.session_state.rewrite_result and not st.session_state.rewrite_generating:
        st.markdown(f"### {lang['rewrite_result_title']}")
        st.text_area("Result", st.session_state.rewrite_result, height=450, key="rewrite_result_preview")
        col_word, col_ppt = st.columns(2)
        with col_word:
            word_buf = generate_word_file(st.session_state.rewrite_result)
            st.download_button(
                label=lang["rewrite_download_word"],
                data=word_buf,
                file_name=lang["rewrite_word_filename"],
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )
        with col_ppt:
            ppt_buf = generate_ppt_file(
                content=st.session_state.rewrite_result,
                title=lang["rewrite_result_title"],
                footer=lang["ppt_footer"],
                end_text=lang["ppt_end_page"]
            )
            st.download_button(
                label=lang["rewrite_download_ppt"],
                data=ppt_buf,
                file_name=lang["rewrite_ppt_filename"],
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )

# ==============================================================================
# 📌 功能6：商务文档翻译（保持原样）
# ==============================================================================
def render_translate():
    st.markdown(lang["translate_tip"])
    col1, col2 = st.columns(2)
    with col1:
        target_lang = st.selectbox(
            lang["translate_target_lang"],
            options=TARGET_LANG_OPTIONS[current_lang],
            index=1
        )
    with col2:
        translate_mode = st.radio(
            lang["translate_mode"],
            options=[lang["translate_mode_text"], lang["translate_mode_file"]],
            horizontal=True
        )
    st.markdown("---")
    source_text = ""
    if translate_mode == lang["translate_mode_text"]:
        source_text = st.text_area(lang["translate_textarea_tip"], height=200)
    else:
        translate_file = st.file_uploader(lang["translate_upload_tip"], type=["txt","docx"])
        if translate_file:
            source_text = read_file(translate_file)
            with st.expander(lang["translate_original_preview"]):
                st.text_area("Original", source_text, height=250)
    st.markdown("---")
    if st.button(lang["translate_btn"], use_container_width=True, disabled=st.session_state.translate_generating):
        if not source_text.strip():
            st.warning(f"{lang['warning']}: {lang['translate_content_empty']}")
        else:
            st.session_state.translate_result = ""
            st.session_state.translate_generating = True
    if st.session_state.translate_generating and not st.session_state.translate_result:
        with st.spinner(lang["translate_loading"]):
            prompt = PROMPT_CONFIG["doc_translate"].format(target_lang=target_lang, text=source_text[:6000])
            result_text = ai_request(prompt, target_lang=st.session_state.language)
            st.session_state.translate_result = result_text
            st.session_state.translate_generating = False
    if st.session_state.translate_result and not st.session_state.translate_generating:
        st.markdown(f"### {lang['translate_result_title']}")
        st.text_area("Result", st.session_state.translate_result, height=400, key="translate_result_preview")
        word_buf = generate_word_file(st.session_state.translate_result)
        st.download_button(
            label=lang["translate_download_btn"],
            data=word_buf,
            file_name=lang["translate_download_filename"],
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            use_container_width=True
        )

# ==============================================================================
# 📌 功能7：无损PDF转Word（纯本地转换，保持原样）
# ==============================================================================
def render_pdf2word():
    st.markdown(lang["pdf2word_tip"])
    pdf_file = st.file_uploader(lang["pdf2word_upload_tip"], type=["pdf"], key="pdf_file")

    if pdf_file:
        with st.spinner(lang["pdf2word_loading"]):
            raw_bytes = pdf_file.read()
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
                f.write(raw_bytes)
                temp_pdf = f.name
            temp_docx = tempfile.mktemp(suffix=".docx")

            try:
                cv = Converter(temp_pdf)
                cv.convert(temp_docx)
                cv.close()

                with open(temp_docx, "rb") as f:
                    word_data = f.read()

                st.success(lang["success"])
                st.download_button(
                    label=lang["pdf2word_download_word"],
                    data=word_data,
                    file_name=lang["pdf2word_word_filename"],
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True
                )
            except Exception as e:
                st.error(f"{lang['error']}: PDF转换失败 - {str(e)}")
            finally:
                try:
                    os.unlink(temp_pdf)
                    os.unlink(temp_docx)
                except:
                    pass

# ==============================================================================
# 🆕 新增功能8：关于本产品 — 产品方法论展示
# ==============================================================================
def render_about():
    st.markdown(f"## {lang['about_title']}")
    st.markdown(f"*{lang['about_subtitle']}*")
    st.markdown("---")
    st.markdown(lang["about_content"])

    # 展示技术指标
    st.markdown("---")
    st.markdown("### 技术指标")
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.metric(label="功能模块", value="11")
    with col2:
        st.metric(label="行业赛道模板", value="6")
    with col3:
        st.metric(label="支持语言", value="5")
    with col4:
        st.metric(label="输出格式", value="Word/PPT")

# ==============================================================================
# 🆕 新增功能9：企业AI就绪度评估
# ==============================================================================
def render_assessment():
    st.markdown(f"## {lang['assessment_title']}")
    st.markdown(f"*{lang['assessment_subtitle']}*")
    st.markdown("---")

    col1, col2 = st.columns(2)
    with col1:
        company_name = st.text_input(lang["assessment_company"], key="assess_company")
        company_size = st.selectbox(
            lang["assessment_size"],
            options=lang["assessment_sizes"],
            key="assess_size"
        )
    with col2:
        industry = st.text_input(lang["assessment_industry"], key="assess_industry")
        ai_level = st.selectbox(
            lang["assessment_existing_ai"],
            options=lang["assessment_ai_levels"],
            key="assess_ai_level"
        )
    revenue = st.text_input(lang["assessment_revenue"], key="assess_revenue", placeholder="如：5000")

    st.markdown("---")
    if st.button(lang["assessment_btn"], use_container_width=True, key="assess_btn"):
        if not company_name.strip():
            st.warning(f"{lang['warning']}: {lang['assessment_company_empty']}")
        else:
            with st.spinner(lang["assessment_loading"]):
                prompt = PROMPT_CONFIG["ai_readiness_assessment"].format(
                    company_name=company_name,
                    industry=industry or "未指定",
                    company_size=company_size,
                    revenue=revenue or "未提供",
                    ai_level=ai_level,
                    compliance_rule=PROMPT_CONFIG["compliance_rule"],
                    formatting_requirement=PROMPT_CONFIG["formatting_requirement"]
                )
                result = ai_request(prompt, target_lang=st.session_state.language)
                st.session_state.assessment_result = result

    if st.session_state.assessment_result:
        st.markdown(f"### {lang['assessment_result_title']}")
        st.markdown(st.session_state.assessment_result)

        # 下载按钮
        col_dl1, col_dl2 = st.columns(2)
        with col_dl1:
            word_buf = generate_word_file(st.session_state.assessment_result)
            st.download_button(
                label="📍 下载Word版评估报告",
                data=word_buf,
                file_name=f"AI就绪度评估报告_{company_name or '企业'}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )
        with col_dl2:
            ppt_buf = generate_ppt_file(
                content=st.session_state.assessment_result,
                title=f"AI就绪度评估 - {company_name or '企业'}",
                footer=lang["ppt_footer"],
                end_text=lang["ppt_end_page"]
            )
            st.download_button(
                label="📊 下载PPT版",
                data=ppt_buf,
                file_name=f"AI就绪度评估报告_{company_name or '企业'}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )

# ==============================================================================
# 🆕 新增功能10：AI合规与溯源 — 交互式
# ==============================================================================
def render_compliance():
    st.markdown(f"## {lang['compliance_title']}")
    st.markdown(f"*{lang['compliance_subtitle']}*")
    st.markdown("---")

    # ── 一、合规检查 ──
    st.markdown(f"### {lang['compliance_check_title']}")
    st.caption(lang["compliance_check_tip"])
    check_input = st.text_area(
        label="compliance_check_area",
        placeholder=lang["compliance_check_placeholder"],
        height=200,
        key="compliance_check_input",
        label_visibility="collapsed"
    )
    if st.button(lang["compliance_check_btn"], use_container_width=True, key="compliance_check_btn"):
        if not check_input.strip():
            st.warning(lang["translate_content_empty"])
        else:
            with st.spinner(lang["compliance_check_loading"]):
                prompt = PROMPT_CONFIG["compliance_check"].format(
                    text_to_check=check_input.strip(),
                    formatting_requirement=PROMPT_CONFIG["formatting_requirement"]
                )
                result = ai_request(prompt, target_lang=st.session_state.language)
                st.session_state.compliance_check_result = result

    if st.session_state.get("compliance_check_result"):
        st.markdown(f"#### {lang['compliance_check_result_title']}")
        st.markdown(st.session_state.compliance_check_result)
        st.markdown("")

    st.markdown("---")

    # ── 二、溯源报告生成 ──
    st.markdown(f"### {lang['compliance_trace_title_new']}")
    st.caption(lang["compliance_trace_tip"])
    trace_input = st.text_input(
        label="compliance_trace_input",
        placeholder=lang["compliance_trace_placeholder"],
        key="compliance_trace_topic"
    )
    if st.button(lang["compliance_trace_btn"], use_container_width=True, key="compliance_trace_btn"):
        if not trace_input.strip():
            st.warning(lang["translate_content_empty"])
        else:
            with st.spinner(lang["compliance_trace_loading"]):
                prompt = PROMPT_CONFIG["traceability_report"].format(
                    report_topic=trace_input.strip(),
                    compliance_rule=PROMPT_CONFIG["compliance_rule"],
                    formatting_requirement=PROMPT_CONFIG["formatting_requirement"]
                )
                result = ai_request(prompt, target_lang=st.session_state.language)
                st.session_state.compliance_trace_result = result

    if st.session_state.get("compliance_trace_result"):
        st.markdown(f"#### {lang['compliance_trace_result_title']}")
        st.markdown(st.session_state.compliance_trace_result)
        col_dl1, col_dl2 = st.columns(2)
        with col_dl1:
            word_buf = generate_word_file(st.session_state.compliance_trace_result)
            st.download_button(
                label=lang["compliance_download_word"],
                data=word_buf,
                file_name="Traceability_Report.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )
        with col_dl2:
            ppt_buf = generate_ppt_file(
                content=st.session_state.compliance_trace_result,
                title="Traceability Report",
                footer=lang["ppt_footer"],
                end_text=lang["ppt_end_page"]
            )
            st.download_button(
                label=lang["compliance_download_ppt"],
                data=ppt_buf,
                file_name="Traceability_Report.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )

    st.markdown("---")

    # ── 三、合规机制参考（静态展示）──
    st.markdown(f"### {lang['compliance_static_title']}")
    st.markdown("")

    # 可信度评分体系
    st.markdown(f"#### {lang['compliance_score_title']}")
    st.markdown(lang["compliance_score_desc"])
    score_data = [
        (lang["compliance_source_gov"], "9-10/10", "🟢", "国家统计局、央行、证监会、国务院部门"),
        (lang["compliance_source_listed"], "8-9/10", "🟢", "上市公司年报、招股说明书、官方公告"),
        (lang["compliance_source_consulting"], "7-8/10", "🟡", "麦肯锡、BCG、贝恩、IDC、Gartner、四大"),
        (lang["compliance_source_media"], "4-6/10", "🟠", "财经媒体、科技媒体、行业自媒体"),
        (lang["compliance_source_unknown"], "1-3/10", "🔴", "未标注来源的数据/信息"),
    ]
    table_md = f"| {lang['compliance_source_levels']} | 评分 | 等级 | 典型来源 |\n|------|------|------|------|\n"
    for row in score_data:
        table_md += f"| {row[0]} | {row[1]} | {row[2]} | {row[3]} |\n"
    st.markdown(table_md)
    st.markdown("")

    # 溯源链路
    st.markdown(f"#### {lang['compliance_trace_title']}")
    st.markdown(lang["compliance_trace_desc"])
    st.info(lang["compliance_trace_flow"])
    st.markdown("")

    # 合规检查开关
    st.markdown(f"#### {lang['compliance_switch_title']}")
    st.markdown(lang["compliance_switch_desc"])
    for item in lang["compliance_switch_items"]:
        st.checkbox(item, value=True, disabled=True, key=f"static_switch_{item[:10]}")
    st.caption("*以上开关在实际使用中可在生成报告前自由切换，当前为默认开启状态。")

# ==============================================================================
# 🆕 新增功能11：战略决策仪表盘 — 交互式艾森豪威尔矩阵
# ==============================================================================
def render_dashboard():
    st.markdown(f"## {lang['dashboard_title']}")
    st.markdown(f"*{lang['dashboard_subtitle']}*")
    st.markdown("---")
    st.info(lang["dashboard_intro"])
    st.markdown("")

    # 用户输入项目清单
    projects_input = st.text_area(
        lang["dashboard_projects_tip"],
        placeholder=lang["dashboard_projects_placeholder"],
        height=180,
        key="dashboard_projects"
    )

    if st.button(lang["dashboard_btn"], use_container_width=True, key="dashboard_gen_btn"):
        if not projects_input.strip():
            st.warning(lang["assessment_company_empty"])
        else:
            with st.spinner(lang["dashboard_loading"]):
                prompt = PROMPT_CONFIG["eisenhower_matrix"].format(
                    projects_input=projects_input.strip(),
                    formatting_requirement=PROMPT_CONFIG["formatting_requirement"]
                )
                result = ai_request(prompt, target_lang=st.session_state.language)
                st.session_state.dashboard_result = result

    if st.session_state.dashboard_result:
        st.markdown(f"### {lang['dashboard_result_title']}")
        st.markdown(st.session_state.dashboard_result)

        col_dl1, col_dl2 = st.columns(2)
        with col_dl1:
            word_buf = generate_word_file(st.session_state.dashboard_result)
            st.download_button(
                label=lang["dashboard_download_word"],
                data=word_buf,
                file_name="Strategic_Decision_Matrix.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )
        with col_dl2:
            ppt_buf = generate_ppt_file(
                content=st.session_state.dashboard_result,
                title="Strategic Decision Matrix",
                footer=lang["ppt_footer"],
                end_text=lang["ppt_end_page"]
            )
            st.download_button(
                label=lang["dashboard_download_ppt"],
                data=ppt_buf,
                file_name="Strategic_Decision_Matrix.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )

# ==============================================================================
# 📋 功能路由映射表 — 7原有 + 4新增
# ==============================================================================
RENDER_FUNC_MAP = {
    "search":       render_search,
    "summary":      render_summary,
    "generate":     render_generate,
    "compare":      render_compare,
    "rewrite":      render_rewrite,
    "translate":    render_translate,
    "pdf2word":     render_pdf2word,
    "about":        render_about,
    "assessment":   render_assessment,
    "compliance":   render_compliance,
    "dashboard":    render_dashboard,
}

# ==============================================================================
# 🎨 Streamlit页面配置与CSS
# ==============================================================================
st.set_page_config(
    page_title=lang["page_title"],
    page_icon=PAGE_CONFIG["page_icon"],
    layout="wide",
    initial_sidebar_state="expanded"
)

st.markdown("""
<style>
#MainMenu {visibility: hidden !important;}
footer {visibility: hidden !important;}
button[title="View fullscreen"] {visibility: hidden !important;}
.stDeployButton {display: none !important;}
::-webkit-scrollbar {display: none !important;}
header {height: 0 !important;background: transparent !important;border: none !important;}
button[aria-label="Open sidebar"] {position: fixed !important;top: 1rem !important;left: 1rem !important;z-index: 9999 !important;background-color: rgba(255,255,255,0.9) !important;border-radius: 50% !important;width: 2.5rem !important;height: 2.5rem !important;box-shadow: 0 2px 8px rgba(0,0,0,0.1) !important;}
button[aria-label="Close sidebar"] {display: none !important;}
.block-container {padding-top: 0 !important;padding-bottom: 1rem !important;max-width: 95% !important;}
h1 {margin-top: 0.5rem !important;}
</style>
""", unsafe_allow_html=True)

st.title(lang["main_title"])
st.markdown("---")

# ==============================================================================
# 🔧 Session State 初始化
# ==============================================================================
if "language" not in st.session_state:
    st.session_state.language = "zh"
if "selected_tab" not in st.session_state:
    st.session_state.selected_tab = ""
if "dashboard_result" not in st.session_state:
    st.session_state.dashboard_result = None
if "compliance_check_result" not in st.session_state:
    st.session_state.compliance_check_result = None
if "compliance_trace_result" not in st.session_state:
    st.session_state.compliance_trace_result = None

# ==============================================================================
# 📂 侧边栏 — 语言切换 + 功能导航
# ==============================================================================
with st.sidebar:
    st.radio(
        lang["lang_select"],
        options=["zh", "en"],
        format_func=lambda x: "中文" if x == "zh" else "English",
        key="language",
        horizontal=True
    )
    st.divider()
    st.header(lang["sidebar_title"])
    if st.session_state.selected_tab == "":
        st.session_state.selected_tab = MENU_LABELS[0]
    st.radio(
        lang["select_func"],
        MENU_LABELS,
        key="selected_tab",
        label_visibility="visible"
    )
    st.markdown("---")
    st.info(lang["sidebar_footer"])

# ==============================================================================
# 🚀 功能路由：根据侧边栏选择渲染对应功能页面
# ==============================================================================
if st.session_state.selected_tab:
    selected_item = MENU_MAP.get(st.session_state.selected_tab)
    if selected_item:
        func_id = selected_item["id"]
        render_func = RENDER_FUNC_MAP.get(func_id)
        if render_func:
            render_func()
